import os
import io
import uuid
import time
import requests
from flask import Flask, request, send_from_directory, url_for, render_template_string, send_file, render_template, Response
from pptx import Presentation
from pptx.util import Inches

import pandas as pd
import geopandas as gpd
import folium
# --- CHANGE: Import MarkerCluster ---
from folium.plugins import MarkerCluster
import googlemaps
from dotenv import load_dotenv

# `unary_union` is deprecated in Shapely 2.1 in favor of `union_all`.  Fall
# back to `unary_union` for older versions so the code runs regardless of the
# installed Shapely release.
try:
    from shapely import union_all
except ImportError:  # pragma: no cover - Shapely < 2.1
    from shapely.ops import unary_union as union_all
from geopy.geocoders import Nominatim
from geopy.extra.rate_limiter import RateLimiter

# --- FIX: Define a base directory to make all file paths absolute ---
basedir = os.path.abspath(os.path.dirname(__file__))

# Load environment variables from .env file
load_dotenv(os.path.join(basedir, '.env'))

# Initialize Flask app and security configs
app = Flask(__name__)
# Debug mode from environment
DEBUG = os.environ.get('DEBUG', 'False').lower() == 'true'
app.config['DEBUG'] = DEBUG
# Maximum upload size (10 MB)
app.config['MAX_CONTENT_LENGTH'] = 10 * 1024 * 1024  # 10MB limit
# Allowed upload MIME types
ALLOWED_MIME_TYPES = [
    'text/csv',
    'application/vnd.ms-excel',
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
]
# CSRF protection
app.config['SECRET_KEY'] = os.environ.get('SECRET_KEY', 'dev-secret-key')
from flask_wtf import CSRFProtect
from werkzeug.utils import secure_filename
from werkzeug.exceptions import RequestEntityTooLarge
csrf = CSRFProtect(app)
# Handle file too large errors
@app.errorhandler(RequestEntityTooLarge)
def handle_file_too_large(e):
    return 'File too large. Maximum size is 10MB.', 413
# Disable template caching for development
app.config['TEMPLATES_AUTO_RELOAD'] = True
app.config['SEND_FILE_MAX_AGE_DEFAULT'] = 0
# Comment out SERVER_NAME for local development
# app.config['SERVER_NAME'] = 'caas-map-old.link-smart-home.com'
app.config['APPLICATION_ROOT'] = '/'
app.config['PREFERRED_URL_SCHEME'] = 'https'
app.config['UPLOAD_FOLDER'] = os.path.join(basedir, 'uploads')

# Add cache-busting headers for remote deployment
@app.after_request
def add_no_cache_headers(response):
    """Add headers to prevent caching in production/remote environments"""
    response.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate, public, max-age=0'
    response.headers['Pragma'] = 'no-cache'
    response.headers['Expires'] = '0'
    return response

# Directories
os.makedirs(os.path.join(basedir, "static", "maps"), exist_ok=True)
os.makedirs(os.path.join(basedir, "static", "img"), exist_ok=True)
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# Preload groups & shapefile
state_groups = pd.read_csv(os.path.join(basedir, "input_csv_files", "group_by_state.csv"))
us_states = gpd.read_file(
    os.path.join(basedir, "us_state_boundary_shapefiles", "ne_10m_admin_1_states_provinces_lakes.shp")
)
us_states = us_states[us_states["admin"] == "United States of America"]
us_states["StateAbbr"] = us_states["iso_3166_2"].str.split("-").str[-1]
us_states = us_states.merge(
    state_groups, left_on="StateAbbr", right_on="State", how="left"
)

GROUP_COLORS = {"Group 1": "#0056b8", "Group 2": "#00a1e0", "Group 3": "#a1d0f3"}

# Google Maps API Key (set your key here or via environment variable)
GOOGLE_MAPS_API_KEY = os.environ.get("GOOGLE_MAPS_API_KEY", "YOUR_GOOGLE_MAPS_API_KEY")

# --- NEW: Helper function to convert hex to rgba for table row coloring ---
def hex_to_rgba(hex_color, alpha=0.2):
    if pd.isna(hex_color):
        return 'rgba(255, 255, 255, 0)' # Return transparent if color is missing
    hex_color = hex_color.lstrip('#')
    r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    return f'rgba({r}, {g}, {b}, {alpha})'

# ------------------------------------------
# Helper function to load an SVG and inject the row's candidate number
# ------------------------------------------
def load_and_inject_svg(svg_path, number_value):
    with open(svg_path, "r", encoding="utf-8") as f:
        svg_content = f.read()
    # Replace the placeholder {{NUMBER}} with the actual numeric value
    svg_content = svg_content.replace("{{NUMBER}}", str(number_value))
    return svg_content


# Pin definitions - move inside route handlers to avoid URL generation issues
PIN_TYPES = {}


def get_pin_types():
    """Get pin types with proper URL generation within request context"""
    return {
        "primary_dark_blue_sphere": {
            "url": url_for("static", filename="img/sphere_pin_primary_dark_blue.svg"),
            "numbered": False,
            "label": "Location Name - Candidates",
        },
        "primary_dark_blue_number": {
            "url": url_for("static", filename="img/number_pin_primary_dark_blue.svg"),
            "numbered": True,
            "label": "Location Name",
        },
        "primary_light_blue_sphere": {
            "url": url_for("static", filename="img/sphere_pin_primary_light_blue.svg"),
            "numbered": False,
            "label": "Location Name - Candidates",
        },
        "primary_light_blue_number": {
            "url": url_for("static", filename="img/number_pin_primary_light_blue.svg"),
            "numbered": True,
            "label": "Location Name",
        },
        "green_sphere": {
            "url": url_for("static", filename="img/sphere_pin_green.svg"),
            "numbered": False,
            "label": "Location Name - Candidates",
        },
        "green_number": {
            "url": url_for("static", filename="img/number_pin_green.svg"),
            "numbered": True,
            "label": "Location Name",
        },
        "secondary_dark_blue_sphere": {
            "url": url_for("static", filename="img/sphere_pin_secondary_dark_blue.svg"),
            "numbered": False,
            "label": "Location Name - Candidates",
        },
        "secondary_dark_blue_number": {
            "url": url_for("static", filename="img/number_pin_secondary_dark_blue.svg"),
            "numbered": True,
            "label": "Location Name",
        },
        "teal_sphere": {
            "url": url_for("static", filename="img/sphere_pin_teal.svg"),
            "numbered": False,
            "label": "Location Name - Candidates",
        },
        "teal_number": {
            "url": url_for("static", filename="img/number_pin_teal.svg"),
            "numbered": True,
            "label": "Location Name",
        },
    }

# Nicer form template with updated styling
UPLOAD_FORM_TEMPLATE = """
<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>Map Uploader</title>
    <style>
      body {
        margin: 0;
        padding: 0;
        font-family: Calibri, sans-serif;
        background: #f4f7fa; /* Soft page background */
      }
      .container {
        max-width: 800px;
        margin: 40px auto;
        background: #fff;
        border-radius: 8px;
        box-shadow: 0 2px 5px rgba(0,0,0,0.15);
        padding: 20px 30px;
      }
      h1 {
        margin-top: 0;
        color: #333;
      }
      label {
        font-weight: bold;
      }
      .pin-selection {
        display: flex;
        flex-wrap: wrap;
        margin-bottom: 20px;
        gap: 10px;
      }
      .pin-option {
        display: inline-flex;
        flex-direction: column;
        align-items: center;
        cursor: pointer;
        transition: transform 0.2s ease;
        border: 2px solid transparent;
        border-radius: 5px;
        padding: 5px;
      }
      .pin-option:hover {
        background: #f0f4f8;
        transform: translateY(-2px);
      }
      .selected-pin {
        border-color: #00a1e0; /* highlight border color when selected */
        background: #e8f6fe;
      }
      .pin-image {
        width: 40px;
        margin-bottom: 5px;
      }
      button {
        background: #0056b8;
        color: #fff;
        border: none;
        border-radius: 4px;
        padding: 10px 16px;
        cursor: pointer;
        font-size: 14px;
      }
      button:hover {
        background: #004494;
      }
      a {
        color: #0056b8;
        text-decoration: none;
        margin-left: 8px;
      }
      a:hover {
        text-decoration: underline;
      }
      .footer-links {
        margin-top: 10px;
      }
      .options {
        margin-top: 20px;
        margin-bottom: 20px;
      }
    </style>
</head>
<body>
<div class="container">
    <h1>Upload CSV/XLS/XLSX to Plot Pins</h1>
    <p><strong>Required columns:</strong> Location Name, ZIP/Postal Code, Electrification Candidates</p>
    <p><strong>Optional columns:</strong> Street Address, City, State</p>
    <p><strong>Accepted file types:</strong> .csv, .xls, .xlsx</p>

        <form action="/assign_pins" method="POST" enctype="multipart/form-data">
            <input type="hidden" name="csrf_token" value="{{ csrf_token() }}"/>
      <label>Choose File:</label><br>
      <input type="file" name="file" accept=".csv,.xls,.xlsx" required/><br><br>

      <div class="options">
        <label>
          <input type="checkbox" name="cluster_pins" value="true" checked> Cluster Pins
        </label>
        <br>
        <label>
          <input type="checkbox" name="show_labels" value="true" checked> Show Labels on Pins
        </label>
      </div>

      <button type="submit">Upload & Generate Map</button>
      <div class="footer-links">
        <a href="/download_template" download="location_pins_template.xlsx">Download Excel Template</a>
      </div>
    </form>
</div>
</body>
</html>
"""

PIN_ASSIGNMENT_TEMPLATE = """
<!DOCTYPE html>
<html>
<head>
    <meta charset=\"UTF-8\">
    <title>Map Uploader - Step 2</title>
    <style>
      body { font-family: Calibri, sans-serif; background: #f4f7fa; margin: 0; padding: 0; }
      .container { max-width: 800px; margin: 40px auto; background: #fff; border-radius: 8px; box-shadow: 0 2px 5px rgba(0,0,0,0.15); padding: 20px 30px; }
      h1 { margin-top: 0; color: #333; }
      label { font-weight: bold; }
      button { background: #0056b8; color: #fff; border: none; border-radius: 4px; padding: 10px 16px; cursor: pointer; font-size: 14px; }
      button:hover { background: #004494; }
      .category-row { display: flex; align-items: center; margin-bottom: 15px; }
      .category-name { width: 250px; font-weight: bold; }
      select { padding: 5px; border-radius: 4px; }
      .pin-preview { width: 40px; height: 40px; margin-left: 15px; object-fit: contain; }
      .pin-style-preview { 
        display: flex; 
        flex-wrap: wrap; 
        align-items: center; 
        gap: 15px; 
        margin-top: 20px; 
        margin-bottom: 20px;
        padding: 15px;
        background-color: #f8f9fa;
        border-radius: 8px;
        border: 1px solid #dee2e6;
      }
      .pin-style-preview > div {
        display: flex;
        flex-direction: column;
        align-items: center;
        text-align: center;
        min-width: 80px;
      }
      .pin-style-preview img {
        width: 30px;
        height: 30px;
        margin-bottom: 5px;
      }
      .pin-style-preview div div {
        font-size: 11px;
        line-height: 1.2;
      }
    </style>
</head>
<body>
<div class=\"container\">
    <h1>Step 2: Assign a Pin to Each Category</h1>
    <p>For each category found in your file, select a pin style from the dropdown.</p>
    <form action=\"/generate_map\" method=\"POST\">
      <input type=\"hidden\" name=\"csrf_token\" value=\"{{ csrf_token() }}\"/>
      <input type=\"hidden\" name=\"filename\" value=\"{{ filename }}\">
      <input type=\"hidden\" name=\"cluster_pins\" value=\"{{ cluster_pins }}\">
      <input type=\"hidden\" name=\"show_labels\" value=\"{{ show_labels }}\">
      {% for category in categories %}
      <div class=\"category-row\">
        <div class=\"category-name\">{{ category }}</div>
        <select name=\"pin_map_{{ category }}\" id=\"select_{{ loop.index }}\" \
                onchange=\"updatePreview({{ loop.index }})\" required>
          {% for pin_key, pin_data in pin_types.items() %}
          <option value=\"{{ pin_key }}\">{{ pin_key.replace('_', ' ')|title }}</option>
          {% endfor %}
        </select>
        <img id=\"preview_{{ loop.index }}\" src=\"\" class=\"pin-preview\">
      </div>
      {% endfor %}
      <div class=\"pin-style-preview\">
        <strong>Pin Style Previews:</strong>
        {% for pin_key, pin_data in pin_types.items() %}
        <div>
          <img src=\"{{ pin_data['url'] }}\" alt=\"{{ pin_key }}\" class=\"pin-preview\">
          <div>{{ pin_key.replace('_', ' ')|title }}</div>
        </div>
        {% endfor %}
      </div>
      <br>
      <button type=\"submit\">Generate Map</button>
    </form>
</div>
<script>
    const pinData = {{ pin_types|tojson }};
    
    function updatePreview(index) {
        const selectEl = document.getElementById('select_' + index);
        const previewEl = document.getElementById('preview_' + index);
        const selectedPinKey = selectEl.value;
        if (pinData[selectedPinKey]) {
            previewEl.src = pinData[selectedPinKey].url;
        }
    }

    document.addEventListener('DOMContentLoaded', function() {
        const selects = document.querySelectorAll('select[id^=\"select_\"]');
        selects.forEach((select, i) => {
            updatePreview(i + 1);
        });
    });
</script>
</body>
</html>
"""


@app.route("/", methods=["GET"])
def upload_form():
    # Render the upload form template
    return render_template("upload_form.html")

@app.route("/assign_pins", methods=["POST"])
def assign_pins():
    try:
        uploaded_file = request.files.get("file")
        if not uploaded_file or not uploaded_file.filename:
            return render_template('error.html', error_code=400, message="No file selected. Please choose a CSV or Excel file."), 400

        # Validate file type and filename
        if uploaded_file.mimetype not in ALLOWED_MIME_TYPES:
            return render_template('error.html', error_code=400, 
                                 message=f"Unsupported file type: {uploaded_file.mimetype}. Please upload a CSV, XLS, or XLSX file."), 400
        
        filename_secure = secure_filename(uploaded_file.filename)
        if not filename_secure:
            return render_template('error.html', error_code=400, message="Invalid filename. Please rename your file."), 400
            
        # Use a unique filename to prevent conflicts
        filename = f"{uuid.uuid4()}_{filename_secure}"
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        
        try:
            uploaded_file.save(filepath)
            app.logger.info(f"File uploaded successfully: {filename}")
        except Exception as e:
            app.logger.error(f"Failed to save uploaded file: {e}")
            return render_template('error.html', error_code=500, message="Failed to save uploaded file. Please try again."), 500

        try:
            if filename and filename.lower().endswith((".xls", ".xlsx")):
                df = pd.read_excel(filepath)
                app.logger.info(f"Successfully read Excel file with {len(df)} rows")
            else:
                df = pd.read_csv(filepath)
                app.logger.info(f"Successfully read CSV file with {len(df)} rows")
        except Exception as e:
            app.logger.error(f"Error reading file {filename}: {e}")
            # Clean up the uploaded file
            try:
                os.remove(filepath)
            except:
                pass
            return render_template('error.html', error_code=400, 
                                 message=f"Error reading file: {str(e)}. Please check your file format and try again."), 400

        required_cols = {"Location Name", "ZIP/Postal Code", "Electrification Candidates", "Category Name"}
        missing_cols = required_cols - set(df.columns)
        if missing_cols:
            # Clean up the uploaded file
            try:
                os.remove(filepath)
            except:
                pass
            return render_template('error.html', error_code=400, 
                                 message=f"Missing required columns: {', '.join(missing_cols)}. Your file must contain: {', '.join(sorted(required_cols))}"), 400

        # Validate data quality
        if len(df) == 0:
            try:
                os.remove(filepath)
            except:
                pass
            return render_template('error.html', error_code=400, message="The uploaded file is empty. Please upload a file with location data."), 400

        df['Category Name'] = df['Category Name'].astype(str).fillna('Uncategorized')
        categories = sorted(df['Category Name'].unique())
        
        app.logger.info(f"File processed successfully: {len(categories)} categories found")

        return render_template(
            "pin_assignment.html",
            filename=filename,
            categories=categories,
            cluster_pins=str(request.form.get("cluster_pins") == "true").lower(),
            show_labels=str(request.form.get("show_labels") == "true").lower()
        )
    except Exception as e:
        app.logger.error(f"An unhandled exception occurred in /assign_pins: {e}", exc_info=True)
        return render_template('error.html', error_code=500, 
                             message="An internal server error occurred. Please try again or contact support if the problem persists."), 500


@app.route("/color_selection", methods=["POST"])
def color_selection():
    """Display the color selection page for state groups."""
    filename = request.form.get("filename")
    cluster_pins = request.form.get("cluster_pins")
    
    if not filename:
        return "Error: No filename provided.", 400
    
    # Get pin assignments from form (now includes type and color)
    pin_assignments = {}
    for key, value in request.form.items():
        if key.startswith("pin_type_"):
            category = key[9:]  # Remove "pin_type_" prefix
            if category not in pin_assignments:
                pin_assignments[category] = {}
            pin_assignments[category]['type'] = value
        elif key.startswith("pin_color_"):
            category = key[10:]  # Remove "pin_color_" prefix
            if category not in pin_assignments:
                pin_assignments[category] = {}
            pin_assignments[category]['color'] = value
    
    # Default colors
    default_colors = {
        "Group 1": "#0056b8",
        "Group 2": "#00a1e0", 
        "Group 3": "#a1d0f3"
    }
    
    return render_template(
        "color_selection.html",
        filename=filename,
        cluster_pins=cluster_pins,
        pin_assignments=pin_assignments,
        default_colors=default_colors
    )


@app.route("/generate_custom_pin_svg")
def generate_custom_pin_svg():
    """Generate a custom pin SVG with specified color and type."""
    pin_type = request.args.get("type", "sphere")  # sphere or number
    color = request.args.get("color", "#0056b8")
    number = request.args.get("number", "1")
    
    if pin_type == "sphere":
        # Generate gradient colors - lighter version for gradient
        base_color = color
        if color.startswith("#"):
            # Create lighter shade for gradient
            r = int(color[1:3], 16)
            g = int(color[3:5], 16) 
            b = int(color[5:7], 16)
            # Make lighter by adding 60 to each component, capped at 255
            r_light = min(255, r + 60)
            g_light = min(255, g + 60)
            b_light = min(255, b + 60)
            light_color = f"#{r_light:02x}{g_light:02x}{b_light:02x}"
        else:
            light_color = color
            
        svg_content = f'''<svg width="50" height="50" viewBox="0 0 75 75" fill="none" xmlns="http://www.w3.org/2000/svg">
  <defs>
    <radialGradient id="pinGradient" cx="30%" cy="30%" r="50%">
      <stop offset="0%" stop-color="{light_color}"/>
      <stop offset="100%" stop-color="{base_color}"/>
    </radialGradient>
  </defs>
  <circle cx="50" cy="30" r="15" fill="url(#pinGradient)" />
  <circle cx="44" cy="24" r="5" fill="#ffffff" fill-opacity="0.3" />
  <path d="M45 45 L50 80" stroke="#888" stroke-width="2" transform="rotate(30, 50, 45)" />
</svg>'''
    else:  # number pin
        # Generate stroke color - lighter version of base color
        base_color = color
        if color.startswith("#"):
            r = int(color[1:3], 16)
            g = int(color[3:5], 16)
            b = int(color[5:7], 16)
            # Make lighter for stroke
            r_stroke = min(255, r + 80)
            g_stroke = min(255, g + 80)
            b_stroke = min(255, b + 80)
            stroke_color = f"#{r_stroke:02x}{g_stroke:02x}{b_stroke:02x}"
        else:
            stroke_color = color
            
        svg_content = f'''<svg xmlns="http://www.w3.org/2000/svg" width="65" height="55" viewBox="0 0 65 55">
  <g transform="translate(0,2)">
  <path fill="{base_color}" stroke="{stroke_color}" stroke-width="2"
        d="M32.5,0
           C22.558,0 14.5,8.058 14.5,18
           c0 13.5 18,32 18,32
           s18-18.5 18-32
           c0-9.942-8.058-18-18-18z" />
  <circle cx="32.5" cy="20" r="10" fill="#ffffff"/>
  <text x="32.5" y="24" text-anchor="middle" font-size="12" font-family="Calibri"
        font-weight="bold" fill="#000000">{number}</text>
  </g>
</svg>'''
    
    return Response(svg_content, mimetype="image/svg+xml")


@app.route("/generate_map", methods=["POST"])
def generate_map():
    filename = request.form.get("filename")
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename or "")
    if not os.path.exists(filepath):
        return "Error: Uploaded file not found. Please try again.", 404

    try:
        if filename and filename.lower().endswith((".xls", ".xlsx")):
            df = pd.read_excel(filepath)
        else:
            df = pd.read_csv(filepath)
    except Exception as e:
        return f"Error reading file: {e}", 400

    df['Category Name'] = df['Category Name'].astype(str).fillna('Uncategorized')

    # Get custom pin assignments (type and color) from form
    pin_assignments = {}
    for key, value in request.form.items():
        if key.startswith("pin_type_"):
            category = key[9:]  # Remove "pin_type_" prefix
            if category not in pin_assignments:
                pin_assignments[category] = {}
            pin_assignments[category]['type'] = value
        elif key.startswith("pin_color_"):
            category = key[10:]  # Remove "pin_color_" prefix
            if category not in pin_assignments:
                pin_assignments[category] = {}
            pin_assignments[category]['color'] = value

    # Get custom state colors from the form, fall back to defaults
    custom_colors = {
        "Group 1": request.form.get("color_group1", "#0056b8"),
        "Group 2": request.form.get("color_group2", "#00a1e0"),
        "Group 3": request.form.get("color_group3", "#a1d0f3")
    }

    cluster_pins = request.form.get("cluster_pins") == 'true'
    show_labels = request.form.get("show_labels", "true") == 'true'  # Default to True

    # Helper function to generate custom pin SVG
    def generate_pin_svg(pin_type, color, number=None):
        if pin_type == "sphere":
            # Generate gradient colors - lighter version for gradient
            if color.startswith("#"):
                r = int(color[1:3], 16)
                g = int(color[3:5], 16) 
                b = int(color[5:7], 16)
                r_light = min(255, r + 60)
                g_light = min(255, g + 60)
                b_light = min(255, b + 60)
                light_color = f"#{r_light:02x}{g_light:02x}{b_light:02x}"
            else:
                light_color = color
                
            return f'''<svg width="50" height="50" viewBox="0 0 75 75" fill="none" xmlns="http://www.w3.org/2000/svg">
  <defs>
    <radialGradient id="pinGradient{hash(color) % 10000}" cx="30%" cy="30%" r="50%">
      <stop offset="0%" stop-color="{light_color}"/>
      <stop offset="100%" stop-color="{color}"/>
    </radialGradient>
  </defs>
  <circle cx="50" cy="30" r="15" fill="url(#pinGradient{hash(color) % 10000})" />
  <circle cx="44" cy="24" r="5" fill="#ffffff" fill-opacity="0.3" />
  <path d="M45 45 L50 80" stroke="#888" stroke-width="2" transform="rotate(30, 50, 45)" />
</svg>'''
        else:  # number pin
            if color.startswith("#"):
                r = int(color[1:3], 16)
                g = int(color[3:5], 16)
                b = int(color[5:7], 16)
                r_stroke = min(255, r + 80)
                g_stroke = min(255, g + 80)
                b_stroke = min(255, b + 80)
                stroke_color = f"#{r_stroke:02x}{g_stroke:02x}{b_stroke:02x}"
            else:
                stroke_color = color
                
            return f'''<svg xmlns="http://www.w3.org/2000/svg" width="65" height="55" viewBox="0 0 65 55">
  <g transform="translate(0,2)">
  <path fill="{color}" stroke="{stroke_color}" stroke-width="2"
        d="M32.5,0
           C22.558,0 14.5,8.058 14.5,18
           c0 13.5 18,32 18,32
           s18-18.5 18-32
           c0-9.942-8.058-18-18-18z" />
  <circle cx="32.5" cy="20" r="10" fill="#ffffff"/>
  <text x="32.5" y="24" text-anchor="middle" font-size="12" font-family="Calibri"
        font-weight="bold" fill="#000000">{number if number else "1"}</text>
  </g>
</svg>'''

    m = folium.Map(
        location=[39.8283, -98.5795],
        zoom_start=5,
        tiles="https://{s}.basemaps.cartocdn.com/light_nolabels/{z}/{x}/{y}{r}.png",
        attr="©OpenStreetMap contributors ©CartoDB",
        zoomSnap=0.01,
        zoomDelta=0.01
    )

    folium.GeoJson(
        data=us_states.__geo_interface__,
        style_function=lambda feat: {
            "fillColor": custom_colors.get(feat["properties"].get("CaaS Group"), "gray"),
            "color": "black",
            "weight": 1,
            "fillOpacity": 1.0,
            "className": ("group1-state" if feat["properties"].get("CaaS Group") == "Group 1" else "")
        },
        tooltip=folium.GeoJsonTooltip(fields=["name"], aliases=["State:"])
    ).add_to(m)

    # Format ZIP codes
    def format_zip(value):
        if pd.isna(value) or value == "":
            return ""
        try:
            return f"{int(float(value)):05d}"
        except (ValueError, TypeError):
            return str(value)
    
    df["ZIP/Postal Code"] = df["ZIP/Postal Code"].apply(format_zip)
    
    # Add missing columns
    for optional_col in ["Street Address", "City", "State"]:
        if optional_col not in df.columns:
            df[optional_col] = ""
        else:
            df[optional_col] = df[optional_col].fillna("")
    
    # Build address strings for geocoding
    def build_address_string(row):
        parts = []
        if row.get("Street Address", "").strip():
            parts.append(row["Street Address"].strip())
        if row.get("City", "").strip():
            parts.append(row["City"].strip())
        if row.get("State", "").strip():
            parts.append(row["State"].strip())
        if row.get("ZIP/Postal Code", "").strip():
            parts.append(row["ZIP/Postal Code"].strip())
        
        if parts:
            return ", ".join(parts) + ", USA"
        else:
            return "USA"

    # Geocode locations using the official Google Maps client
    lat_list, lon_list = [], []
    api_key = GOOGLE_MAPS_API_KEY
    # print(f"Using Google Maps API key: {api_key}")  # Debug print
    gmaps = googlemaps.Client(key=api_key) # Initialize the client

    # Geocode locations with improved fallback and logging
    geocoding_stats = {"full_address": 0, "zip_only": 0, "state_centroid": 0, "failed": 0}
    
    for counter, (idx, row) in enumerate(df.iterrows(), 1):
        lat, lon = None, None # Reset lat/lon for each row
        location_name = str(row.get("Location Name", f"Location {counter}"))
        addr_str = build_address_string(row)
        geocoding_method = "failed"
        
        try:
            # Primary: Try full address geocoding
            geocode_result = gmaps.geocode(addr_str)
            if geocode_result:
                loc = geocode_result[0]['geometry']['location']
                lat, lon = loc['lat'], loc['lng']
                geocoding_method = "full_address"
                geocoding_stats["full_address"] += 1

            # Fallback 1: ZIP only if the full address failed
            if lat is None and row.get("ZIP/Postal Code", "").strip():
                zip_addr = f"{row['ZIP/Postal Code']}, USA"
                geocode_result = gmaps.geocode(zip_addr)
                if geocode_result:
                    loc = geocode_result[0]['geometry']['location']
                    lat, lon = loc['lat'], loc['lng']
                    geocoding_method = "zip_only"
                    geocoding_stats["zip_only"] += 1
                    app.logger.info(f"Used ZIP fallback for {location_name}")

            # Fallback 2: State centroid
            if lat is None and row.get("State", "").strip():
                state_abbr = row["State"].strip()
                if state_abbr in us_states["StateAbbr"].values:
                    state_geom = us_states[us_states["StateAbbr"] == state_abbr].geometry.iloc[0]
                    centroid = state_geom.centroid
                    lat, lon = centroid.y, centroid.x
                    geocoding_method = "state_centroid"
                    geocoding_stats["state_centroid"] += 1
                    app.logger.warning(f"Used state centroid fallback for {location_name} in {state_abbr}")
            
            # Final check for complete failure
            if lat is None or lon is None:
                geocoding_stats["failed"] += 1
                app.logger.error(f"Complete geocoding failure for {location_name} with address: {addr_str}")

        except Exception as e:
            geocoding_stats["failed"] += 1
            app.logger.error(f"Geocoding exception for {location_name}: {e}")
            pass

        lat_list.append(lat)
        lon_list.append(lon)

        # The time.sleep(0.05) is NO LONGER NEEDED. The client handles rate limits.

    df['Latitude'] = lat_list
    df['Longitude'] = lon_list
    
    # Log detailed geocoding statistics
    total_locations = len(df)
    successful_geocoding = sum(1 for lat in lat_list if lat is not None)
    app.logger.info(f"Geocoding complete: {successful_geocoding}/{total_locations} locations processed")
    app.logger.info(f"Geocoding breakdown - Full address: {geocoding_stats['full_address']}, ZIP only: {geocoding_stats['zip_only']}, State centroid: {geocoding_stats['state_centroid']}, Failed: {geocoding_stats['failed']}")
    
    if geocoding_stats["failed"] > 0:
        app.logger.warning(f"{geocoding_stats['failed']} locations could not be geocoded and will not appear on the map")
    
    # Log geocoding results (legacy compatibility)
    failed_geocoding = [lat for lat in lat_list if lat is None]
    if failed_geocoding:
        app.logger.warning(f"Failed to geocode {len(failed_geocoding)} addresses")
    
    successful_geocoding = sum(1 for lat in lat_list if lat is not None)
    app.logger.info(f"Successfully geocoded {successful_geocoding}/{len(df)} locations")

    df = df.merge(us_states[['StateAbbr', 'CaaS Group']], left_on='State', right_on='StateAbbr', how='left')

    if cluster_pins:
        icon_create_function = """
        function(cluster) {
            var childCount = cluster.getChildCount();
            var color = '#6bc04b'; // green default
            if (childCount >= 100) {
                color = '#0056b8'; // dark blue
            } else if (childCount >= 50) {
                color = '#00a1e0'; // light blue
            } else if (childCount >= 10) {
                color = '#00bfae'; // teal
            }
            return new L.DivIcon({
                html: '<div style="background:' + color + '"><span>' + childCount + '</span></div>',
                className: 'marker-cluster',
                iconSize: new L.Point(40, 40)
            });
        }
        """
        marker_layer = MarkerCluster(icon_create_function=icon_create_function).add_to(m)
    else:
        marker_layer = m

    legend_items = {}
    table_rows = []

    for _, row in df.iterrows():
        lat, lon = row["Latitude"], row["Longitude"]
        if pd.notnull(lat) and pd.notnull(lon):
            category = row['Category Name']
            
            # Get pin assignment for this category
            pin_assignment = pin_assignments.get(category, {
                'type': 'sphere', 
                'color': '#00a1e0'
            })
            pin_type = pin_assignment.get('type', 'sphere')
            pin_color = pin_assignment.get('color', '#00a1e0')

            # Add to legend (generate a sample SVG for legend)
            if category not in legend_items:
                legend_items[category] = f"/generate_custom_pin_svg?type={pin_type}&color={pin_color.replace('#', '%23')}"

            # Generate custom pin SVG with conditional labels
            if pin_type == 'number':
                pin_svg = generate_pin_svg(pin_type, pin_color, row['Electrification Candidates'])
                if show_labels:
                    html = (
                        f'<div class="div-icon-container">'
                        f'<div class="pin-image-wrapper">{pin_svg}</div>'
                        f'<div class="custom-label-text">{row["Location Name"]}</div>'
                        f'</div>'
                    )
                    icon = folium.DivIcon(html=html, icon_size=(150, 110), icon_anchor=(32, 80))
                else:
                    html = (
                        f'<div class="div-icon-container">'
                        f'<div class="pin-image-wrapper">{pin_svg}</div>'
                        f'</div>'
                    )
                    icon = folium.DivIcon(html=html, icon_size=(65, 55), icon_anchor=(32, 55))
            else:  # sphere
                pin_svg = generate_pin_svg(pin_type, pin_color)
                if show_labels:
                    html = (
                        f'<div class="div-icon-container">'
                        f'<div class="pin-image-wrapper">{pin_svg}</div>'
                        f'<div class="custom-label-text">{row["Location Name"]}</div>'
                        f'</div>'
                    )
                    icon = folium.DivIcon(html=html, icon_size=(150, 80), icon_anchor=(25, 50))
                else:
                    html = (
                        f'<div class="div-icon-container">'
                        f'<div class="pin-image-wrapper">{pin_svg}</div>'
                        f'</div>'
                    )
                    icon = folium.DivIcon(html=html, icon_size=(50, 50), icon_anchor=(25, 50))

            # Create marker with popup information
            popup_text = f"<strong>{row['Location Name']}</strong><br>Electrification Candidates: {row['Electrification Candidates']}"
            folium.Marker(
                location=[lat, lon], 
                icon=icon,
                popup=folium.Popup(popup_text, max_width=200)
            ).add_to(marker_layer)

        n = row.get('Electrification Candidates', 1)
        try:
            n = int(n)
        except Exception:
            n = 1
        if n >= 100:
            color = '#0056b8'
        elif n >= 50:
            color = '#00a1e0'
        elif n >= 10:
            color = '#00bfae'
        elif n >= 2:
            color = '#6bc04b'
        else:
            color = None
        if color:
            rgba_color = hex_to_rgba(color, alpha=0.25)
        else:
            rgba_color = "transparent"
        table_rows.append({
            "location": row['Location Name'],
            "candidates": row['Electrification Candidates'],
            "color": rgba_color
        })

    map_html = m._repr_html_()

    # Save the map to a file and redirect to intermediate page
    map_id = str(uuid.uuid4())
    map_path = os.path.join(basedir, "static", "maps", f"{map_id}.html")
    
    # Render the map template to a file
    map_content = render_template("map_template.html", 
                                  map_html=map_html, 
                                  legend_items=legend_items, 
                                  table_rows=table_rows,
                                  group_colors=custom_colors)
    
    with open(map_path, 'w', encoding='utf-8') as f:
        f.write(map_content)
    
    # Clean up: Delete the uploaded file after successful map generation
    try:
        if os.path.exists(filepath):
            os.remove(filepath)
            app.logger.info(f"Successfully deleted uploaded file: {filepath}")
        else:
            app.logger.warning(f"Uploaded file not found for deletion: {filepath}")
    except PermissionError as e:
        app.logger.warning(f"Permission denied when deleting uploaded file {filepath}: {e}")
    except FileNotFoundError as e:
        app.logger.info(f"Uploaded file already deleted or not found: {filepath}")
    except Exception as e:
        app.logger.error(f"Unexpected error deleting uploaded file {filepath}: {e}")
    
    # Redirect to intermediate page with Start Over button
    # Render the map generation success page with geocoding statistics
    return render_template("map_success.html", 
                          map_id=map_id,
                          geocoding_stats={
                              'total': total_locations,
                              'successful': successful_geocoding,
                              'failed': geocoding_stats['failed']
                          })


@app.route("/map/<map_id>")
def serve_map(map_id):
    return send_from_directory(os.path.join(basedir, "static", "maps"), f"{map_id}.html")


@app.route("/ppt/<map_id>")
def download_ppt(map_id):
    maps_dir = os.path.join(basedir, "static", "maps")
    html_path = os.path.join(maps_dir, f"{map_id}.html")
    if not os.path.isfile(html_path):
        return "Error: Map not found.", 404

    link = url_for("serve_map", map_id=map_id, _external=True)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    if title:
        title.text = "Interactive Map"

    try:
        slide.shapes.add_ole_object(
            html_path,
            prog_id="htmlfile",
            left=Inches(0.5),
            top=Inches(1.5),
            width=Inches(9),
            height=Inches(5),
        )
    except Exception:
        box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
        tf = box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Open Map"
        run.hyperlink.address = link

    ppt_filename = f"{map_id}.pptx"
    ppt_path = os.path.join(maps_dir, ppt_filename)
    prs.save(ppt_path)
    return send_from_directory(maps_dir, ppt_filename, as_attachment=True)


@app.route("/download_template")
def download_template():
    df = pd.DataFrame({
        "Location Name": ["Example A", "Example B"],
        "Street Address": ["123 Main St", ""],
        "City": ["Anytown", ""],
        "State": ["CA", "TX"],
        "ZIP/Postal Code": ["12345", "67890"],
        "Electrification Candidates": [10, 5],
        "Category Name": ["Retail", "Warehouse"],
    })
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        df.to_excel(writer, index=False, sheet_name='locations')
    output.seek(0)
    return send_file(output, download_name="location_pins_template.xlsx", as_attachment=True)


@app.route("/map_output/<map_id>")
def map_output(map_id):
    map_link = url_for("serve_map", map_id=map_id, _external=True)
    return f"""
    <!DOCTYPE html>
    <html>
    <head>
        <title>Map Output</title>
    </head>
    <body>
        <h1>Your Map is Ready</h1>
        <p><a href=\"{map_link}\" target=\"_blank\">Click here to view your map</a></p>
        <p><a href=\"/\">Start Over</a></p>
    </body>
    </html>
    """

# ------------------------------------------
# Error handling framework
# ------------------------------------------
@app.errorhandler(400)
def bad_request(e):
    return render_template('error.html', error_code=400, message=str(e)), 400

@app.errorhandler(404)
def not_found(e):
    return render_template('error.html', error_code=404, message='Resource not found'), 404

@app.errorhandler(500)
def internal_error(e):
    return render_template('error.html', error_code=500, message='An internal server error occurred'), 500


if __name__ == "__main__":
    app.run(host='0.0.0.0', port=5050, debug=app.config.get('DEBUG', False))
